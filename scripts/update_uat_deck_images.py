from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / 'docs' / 'uat' / 'uat_overview.pptx'
SCREENSHOT_DIR = (
    ROOT
    / 'frontend'
    / 'cypress'
    / 'screenshots'
    / 'uat_capture_pages.spec.js'
)

SLIDE_DEFINITIONS: list[dict[str, object]] = [
    {
        'slide': 4,
        'title': 'Login Page',
        'image': 'uat_login_page.png',
        'header': 'Login establishes authenticated access and league context for every workflow.',
        'purpose': 'Allow users to authenticate quickly and enter the app with the correct league selected.',
        'bullets': [
            'User enters credentials and league details from a single entry screen.',
            'Validation errors are clear and actionable before retry.',
            'Successful sign-in routes directly into the authenticated app shell.',
        ],
    },
    {
        'slide': 5,
        'title': 'Home Page',
        'image': 'uat_home_page.png',
        'header': 'Home is the launch point for standings, status, and fast navigation to core tools.',
        'purpose': 'Give managers immediate league context and quick access to team and roster decisions.',
        'bullets': [
            'Surface standing, activity, and key league state at a glance.',
            'Provide clear navigation to roster, waivers, and commissioner workflows.',
            'Keep summary content readable without requiring scrolling for first actions.',
        ],
    },
    {
        'slide': 6,
        'title': 'War Room',
        'image': 'uat_war_room_page.png',
        'header': 'War Room supports live draft execution with board state and nomination controls.',
        'purpose': 'Enable managers and commissioners to execute draft actions with full context.',
        'bullets': [
            'Draft board shows picks, owners, and roster pressure in real time.',
            'Nomination and bid flow keeps amounts and owner turns visible.',
            'Best available panel augments pick quality during draft windows.',
        ],
    },
    {
        'slide': 7,
        'title': 'Chat Advisor',
        'image': 'uat_chat_advisor_page.png',
        'header': 'Chat Advisor provides conversational guidance for lineup and roster choices.',
        'purpose': 'Help users ask natural-language questions and receive actionable football guidance.',
        'bullets': [
            'Prompt entry and response pane remain readable during long exchanges.',
            'Response format supports quick decisions before lineup lock.',
            'Error and loading states are explicit when advisor data is unavailable.',
        ],
    },
    {
        'slide': 8,
        'title': 'Draft Day Analyzer',
        'image': 'uat_draft_day_analyzer_page.png',
        'header': 'Draft Day Analyzer evaluates value and roster construction signals during draft.',
        'purpose': 'Identify high-value picks and positional risk before submitting bids or selections.',
        'bullets': [
            'Show value indicators for candidate players by position.',
            'Highlight roster balance and positional scarcity impact.',
            'Support pick strategy without requiring manual spreadsheet analysis.',
        ],
    },
    {
        'slide': 9,
        'title': 'My Team',
        'image': 'uat_my_team_page.png',
        'header': 'My Team is the primary roster management surface for weekly lineup control.',
        'purpose': 'Let users review roster slots, submit legal lineups, and evaluate bench options.',
        'bullets': [
            'Starter and bench groupings are clearly separated by role.',
            'Lineup actions enforce league rules with clear validation feedback.',
            'Trade and player detail actions are accessible from the roster cards.',
        ],
    },
    {
        'slide': 10,
        'title': 'Matchups',
        'image': 'uat_matchups_page.png',
        'header': 'Matchups summarizes head-to-head outcomes and weekly scoreboard context.',
        'purpose': 'Allow users to compare teams quickly and navigate to game-level detail.',
        'bullets': [
            'Weekly matchup cards are visible with team totals and status.',
            'Navigation between matchup weeks is consistent and predictable.',
            'Entry points to Game Center are obvious for deeper analysis.',
        ],
    },
    {
        'slide': 11,
        'title': 'Game Center',
        'image': 'uat_game_center_page.png',
        'header': 'Game Center provides detailed scoring context for a single matchup.',
        'purpose': 'Help users understand where points are gained or lost across starters and bench.',
        'bullets': [
            'Per-player scoring context is visible for both teams.',
            'Remaining game and matchup trend information stays readable.',
            'Comparative view supports start/sit and waiver decisions.',
        ],
    },
    {
        'slide': 12,
        'title': 'Waiver Wire',
        'image': 'uat_waiver_wire_page.png',
        'header': 'Waiver Wire supports claim, drop, and priority-based player acquisition.',
        'purpose': 'Let managers add players safely while preserving waiver policy constraints.',
        'bullets': [
            'Player pool is filterable and claim actions are clearly presented.',
            'Claim confirmation protects users from accidental submissions.',
            'Roster-full flow clearly transitions into drop-player selection.',
        ],
    },
    {
        'slide': 13,
        'title': 'Keepers',
        'image': 'uat_keepers_page.png',
        'header': 'Keepers allows managers to select and lock eligible retained players.',
        'purpose': 'Support off-season retention with rule-aware selections and lock workflow.',
        'bullets': [
            'Eligible players and keeper cost context are visible before submit.',
            'Submit and lock actions are separated to prevent premature finalization.',
            'Keeper status remains visible after save and refresh cycles.',
        ],
    },
    {
        'slide': 14,
        'title': 'Analytics',
        'image': 'uat_analytics_page.png',
        'header': 'Analytics Dashboard is the chart navigation hub for league insight modules.',
        'purpose': 'Provide one-click access to trend, value, efficiency, and trade analysis tools.',
        'bullets': [
            'Chart selector buttons expose each analytics module explicitly.',
            'Selected chart panel renders within a consistent card surface.',
            'Users can pivot between chart types without leaving the page.',
        ],
    },
    {
        'slide': 15,
        'title': 'Playoff Bracket',
        'image': 'uat_playoff_bracket_page.png',
        'header': 'Playoff Bracket visualizes elimination paths and round progression.',
        'purpose': 'Enable users to verify playoff seeding and advancement paths at a glance.',
        'bullets': [
            'Round and matchup progression is structurally easy to follow.',
            'Team placement reflects latest playoff state from league data.',
            'Bracket remains readable on standard desktop viewport captures.',
        ],
    },
    {
        'slide': 16,
        'title': 'Commissioner Dashboard',
        'image': 'uat_commissioner_dashboard_page.png',
        'header': 'Commissioner Dashboard centralizes league governance and admin entry points.',
        'purpose': 'Give commissioners a single launch surface for all league configuration actions.',
        'bullets': [
            'Control cards route to lineup, waiver, trade, and division management.',
            'Commissioner-only access boundaries are visible and testable.',
            'Primary edit actions are clearly labeled per governance area.',
        ],
    },
    {
        'slide': 17,
        'title': 'Commissioner - Manage Owners',
        'image': 'uat_commissioner_manage_owners_page.png',
        'header': 'Manage Owners controls league member profile and team metadata.',
        'purpose': 'Allow commissioner updates to owner records with safe edit/save flow.',
        'bullets': [
            'Owner rows are editable with clear validation boundaries.',
            'Unsaved and saved states are distinct for change confidence.',
            'Updated owner metadata persists and is visible after refresh.',
        ],
    },
    {
        'slide': 18,
        'title': 'Commissioner - Lineup Rules',
        'image': 'uat_commissioner_lineup_rules_page.png',
        'header': 'Lineup Rules configures legal roster slot behavior for all teams.',
        'purpose': 'Define and enforce lineup constraints that downstream team pages must honor.',
        'bullets': [
            'Rule controls are editable with explicit save action.',
            'Lineup validation behavior maps to configured slot rules.',
            'Rule changes are durable and visible after navigation.',
        ],
    },
    {
        'slide': 19,
        'title': 'Commissioner - Waiver Rules',
        'image': 'uat_commissioner_waiver_rules_page.png',
        'header': 'Waiver Rules configures timing, priority, and claim policy behavior.',
        'purpose': 'Allow commissioners to update waiver policy that directly governs claim outcomes.',
        'bullets': [
            'Current and editable waiver settings are shown together for review.',
            'Update Waiver Rules action writes policy changes explicitly.',
            'Saved settings persist and reflect on member-facing waiver flows.',
        ],
    },
    {
        'slide': 20,
        'title': 'Commissioner - Manage Trades',
        'image': 'uat_commissioner_manage_trades_page.png',
        'header': 'Manage Trades is the commissioner approval/rejection workflow for trade requests.',
        'purpose': 'Support auditable trade moderation with immediate status transitions.',
        'bullets': [
            'Pending trades are visible with ownership context.',
            'Approve/reject actions are explicit and actionable.',
            'Trade status updates are durable after refresh.',
        ],
    },
    {
        'slide': 21,
        'title': 'Commissioner - Manage Divisions',
        'image': 'uat_commissioner_manage_divisions_page.png',
        'header': 'Manage Divisions defines group structure used by standings and scheduling.',
        'purpose': 'Allow commissioner to edit division assignments without ambiguity.',
        'bullets': [
            'Division assignment controls are visible and deterministic.',
            'Save/confirm path protects against partial updates.',
            'Division changes propagate into standings context.',
        ],
    },
    {
        'slide': 22,
        'title': 'Admin Settings',
        'image': 'uat_admin_settings_page.png',
        'header': 'Admin Settings manages global controls beyond league-level permissions.',
        'purpose': 'Provide site administrators with role and platform-level control points.',
        'bullets': [
            'Admin-only controls are isolated from standard user pages.',
            'Role management actions include clear confirmation feedback.',
            'Operational updates remain visible after refresh.',
        ],
    },
    {
        'slide': 23,
        'title': 'Report a Bug + UAT Handoff',
        'image': 'uat_bug_report_page.png',
        'header': 'Bug report flow closes the loop between UAT findings and engineering action.',
        'purpose': 'Ensure tester evidence and defect details are captured before release decisions.',
        'bullets': [
            'Bug report form captures reproducible issue context for triage.',
            'UAT workbook references are visible for pass/fail traceability.',
            'Handoff path from tester to fix owner is explicit.',
        ],
    },
    {
        'slide': 24,
        'title': 'Key Modals and Overlays',
        'image': 'uat_commissioner_draft_budgets_modal.png',
        'header': 'Modal states are captured separately for workflows requiring explicit confirmation.',
        'purpose': 'Highlight commissioner-only overlay behavior for draft budget governance.',
        'bullets': [
            'Set Draft Budgets modal opens from Commissioner Dashboard.',
            'Budget entry controls and action placement are visible and readable.',
            'Modal close and save affordances are obvious for testers.',
        ],
    },
    {
        'slide': 25,
        'title': 'Trade Proposal Modal',
        'image': 'uat_trade_proposal_modal.png',
        'header': 'Modal states are captured separately for workflows requiring explicit confirmation.',
        'purpose': 'Capture trade proposal inputs before submit so field behavior can be validated.',
        'bullets': [
            'Manager, offered player, and requested player fields are visible.',
            'Submit and cancel controls are accessible and unambiguous.',
            'Labels match the user-facing trade workflow language.',
        ],
    },
    {
        'slide': 26,
        'title': 'Player Season Performance Modal',
        'image': 'uat_player_season_performance_modal.png',
        'header': 'Modal states are captured separately for workflows requiring explicit confirmation.',
        'purpose': 'Display player season detail context used for start/sit and trade decisions.',
        'bullets': [
            'Identity block and season metrics render in one consolidated modal.',
            'Key stats remain legible on desktop viewport capture.',
            'Exit control closes modal without losing page state.',
        ],
    },
    {
        'slide': 27,
        'title': 'Waiver Wire Modal Targets',
        'image': 'uat_waiver_confirm_modal.png',
        'header': 'Modal states are captured separately for workflows requiring explicit confirmation.',
        'purpose': 'Validate waiver claim confirmation and roster-full drop workflow behavior.',
        'bullets': [
            'Confirm Waiver Action modal is captured when claim controls are available.',
            'Roster Full drop-player modal is captured for forced replacement flow.',
            'Deterministic fallback artifacts are generated if claim controls are suppressed.',
        ],
    },
    {
        'slide': 28,
        'title': 'Commissioner - Keeper Rules',
        'image': 'uat_commissioner_keeper_rules_page.png',
        'header': 'Keeper Rules control retention limits and commissioner override operations.',
        'purpose': 'Allow commissioner to configure keeper policy and admin-level keeper operations.',
        'bullets': [
            'Max keepers, max years, and deadlines are editable from one form.',
            'Save action applies keeper settings for downstream owner selection.',
            'Admin operations (veto/reset) are visible for governance workflows.',
        ],
    },
    {
        'slide': 29,
        'title': 'Commissioner - General Ledger Statement',
        'image': 'uat_commissioner_ledger_statement_page.png',
        'header': 'Ledger Statement provides auditable transaction history by selected owner.',
        'purpose': 'Enable commissioner audit of financial and budget-related transactions.',
        'bullets': [
            'Owner, currency, and season filters constrain statement output.',
            'Summary cards expose balance and entry counts before table review.',
            'Entry table includes type, direction, amount, and references.',
        ],
    },
    {
        'slide': 30,
        'title': 'Owner - My Ledger Statement',
        'image': 'uat_owner_ledger_statement_page.png',
        'header': 'Owner ledger view exposes personal transaction history and running balance.',
        'purpose': 'Allow each manager to self-audit credits, debits, and transaction notes.',
        'bullets': [
            'Owner can filter by currency and season year.',
            'Balance is visible without leaving the statement page.',
            'Transaction table remains readable with references and notes.',
        ],
    },
    {
        'slide': 31,
        'title': 'War Room - Show Best Available',
        'image': 'uat_war_room_best_available_panel.png',
        'header': 'Best Available panel surfaces filtered candidate targets during active drafting.',
        'purpose': 'Support faster draft decisions with sortable player and value context.',
        'bullets': [
            'Panel includes rank, player identity, and projected value fields.',
            'Position filters narrow candidate list by roster need.',
            'Sidebar can be toggled without losing draft board visibility.',
        ],
    },
    {
        'slide': 32,
        'title': 'Player Card Detail',
        'image': 'uat_player_identity_card_modal.png',
        'header': 'Player card detail highlights identity and context before deeper stat review.',
        'purpose': 'Ensure player identity data is clear at the top of detail modal workflows.',
        'bullets': [
            'Player name, position, and NFL team are immediately visible.',
            'Identity section anchors the season metrics shown below.',
            'Card layout remains readable for quick in-draft decisions.',
        ],
    },
    {
        'slide': 33,
        'title': 'Analytics - Draft Value Analysis',
        'image': 'uat_analytics_draft_value_page.png',
        'header': 'Draft Value Analysis compares pick outcomes and value efficiency.',
        'purpose': 'Help users identify where draft capital produced surplus or deficit value.',
        'bullets': [
            'Draft Value Analysis module is selected from analytics controls.',
            'Chart panel reflects selected analytics mode in-place.',
            'Outputs support trade and roster planning decisions.',
        ],
    },
    {
        'slide': 34,
        'title': 'Analytics - Efficiency Leaderboard',
        'image': 'uat_analytics_efficiency_page.png',
        'header': 'Efficiency Leaderboard ranks managers by performance and execution quality.',
        'purpose': 'Expose comparative manager efficiency for weekly strategic adjustments.',
        'bullets': [
            'Leaderboard mode is selected from analytics module buttons.',
            'Manager ranking output is readable within analytics panel.',
            'Comparative results support lineup and waiver prioritization.',
        ],
    },
    {
        'slide': 35,
        'title': 'Analytics - Trade Analyzer',
        'image': 'uat_analytics_trade_analyzer_page.png',
        'header': 'Trade Analyzer models likely impact before proposing a trade package.',
        'purpose': 'Allow users to evaluate trade scenarios before submitting proposals.',
        'bullets': [
            'Trade Analyzer mode is selected from analytics controls.',
            'Scenario output appears within the standard analytics card layout.',
            'Result context helps users avoid negative-value trade submissions.',
        ],
    },
]

# Fallback image box from the template dimensions.
DEFAULT_LEFT = Inches(0.75)
DEFAULT_TOP = Inches(1.49)
DEFAULT_WIDTH = Inches(7.82)
DEFAULT_HEIGHT = Inches(5.24)

# Template text geometry
DECK_X, DECK_Y, DECK_W, DECK_H = Inches(0.6), Inches(0.12), Inches(7.2), Inches(0.28)
TITLE_X, TITLE_Y, TITLE_W, TITLE_H = Inches(0.6), Inches(0.45), Inches(7.4), Inches(0.5)
HEADER_X, HEADER_Y, HEADER_W, HEADER_H = Inches(0.6), Inches(1.00), Inches(12.2), Inches(0.4)
PURPOSE_HEAD_X, PURPOSE_HEAD_Y, PURPOSE_HEAD_W, PURPOSE_HEAD_H = Inches(8.72), Inches(1.54), Inches(4.29), Inches(0.39)
PURPOSE_BODY_X, PURPOSE_BODY_Y, PURPOSE_BODY_W, PURPOSE_BODY_H = Inches(8.95), Inches(2.02), Inches(3.95), Inches(1.45)
CORE_HEAD_X, CORE_HEAD_Y, CORE_HEAD_W, CORE_HEAD_H = Inches(8.72), Inches(3.82), Inches(4.29), Inches(0.39)
CORE_BODY_X, CORE_BODY_Y, CORE_BODY_W, CORE_BODY_H = Inches(8.95), Inches(4.28), Inches(3.95), Inches(2.05)

BLACK = RGBColor(0, 0, 0)
BORDER = RGBColor(18, 72, 98)


def find_or_make_picture_box(slide):
    """Return picture placeholder geometry; preserve existing picture position when present."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return shape.left, shape.top, shape.width, shape.height, shape
    return DEFAULT_LEFT, DEFAULT_TOP, DEFAULT_WIDTH, DEFAULT_HEIGHT, None


def replace_slide_image(slide, image_path: Path) -> None:
    left, top, width, height, old_picture = find_or_make_picture_box(slide)

    if old_picture is not None:
        element = old_picture._element
        element.getparent().remove(element)

    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def ensure_slide_count(prs: Presentation, slide_number: int) -> None:
    blank_layout = prs.slide_layouts[6]
    while len(prs.slides) < slide_number:
        prs.slides.add_slide(blank_layout)


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        remove_shape(shape)


def set_single_line_text(shape, text: str, size: int, bold: bool = False) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    if p.runs:
        run = p.runs[0]
        run.font.name = 'Aptos'
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = BLACK


def set_wrapped_body_text(shape, text: str, size: int = 14) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)

    p = tf.paragraphs[0]
    p.text = text
    p.level = 0
    if p.runs:
        run = p.runs[0]
        run.font.name = 'Aptos'
        run.font.size = Pt(size)
        run.font.bold = False
        run.font.color.rgb = BLACK


def force_bullet(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    bu_char = OxmlElement('a:buChar')
    bu_char.set('char', '•')
    p_pr.append(bu_char)


def set_bullet_lines(shape, lines: list[str], size: int = 14) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)

    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        force_bullet(p)
        if p.runs:
            run = p.runs[0]
            run.font.name = 'Aptos'
            run.font.size = Pt(size)
            run.font.bold = False
            run.font.color.rgb = BLACK


def add_labeled_box(slide, left, top, width, height, label: str):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = BORDER
    set_single_line_text(box, label, 14, bold=False)
    return box


def upsert_feature_slide(
    prs: Presentation,
    *,
    slide_number: int,
    title: str,
    screenshot_path: Path,
    header_text: str,
    purpose_text: str,
    bullets: list[str],
) -> bool:
    ensure_slide_count(prs, slide_number)
    slide = prs.slides[slide_number - 1]
    clear_slide(slide)

    # Left image area
    if screenshot_path.exists():
        slide.shapes.add_picture(
            str(screenshot_path),
            DEFAULT_LEFT,
            DEFAULT_TOP,
            width=DEFAULT_WIDTH,
            height=DEFAULT_HEIGHT,
        )
        image_added = True
    else:
        image_added = False

    # Text structure matches other walkthrough slides
    deck_shape = slide.shapes.add_textbox(DECK_X, DECK_Y, DECK_W, DECK_H)
    set_single_line_text(deck_shape, 'PPL Insight Hub - UAT Overview', 18, bold=False)

    title_shape = slide.shapes.add_textbox(TITLE_X, TITLE_Y, TITLE_W, TITLE_H)
    set_single_line_text(title_shape, title, 28, bold=True)

    header_shape = slide.shapes.add_textbox(HEADER_X, HEADER_Y, HEADER_W, HEADER_H)
    set_single_line_text(header_shape, header_text, 18, bold=False)

    add_labeled_box(slide, PURPOSE_HEAD_X, PURPOSE_HEAD_Y, PURPOSE_HEAD_W, PURPOSE_HEAD_H, 'PURPOSE')
    purpose_body = slide.shapes.add_textbox(PURPOSE_BODY_X, PURPOSE_BODY_Y, PURPOSE_BODY_W, PURPOSE_BODY_H)
    set_wrapped_body_text(purpose_body, purpose_text, 14)

    add_labeled_box(slide, CORE_HEAD_X, CORE_HEAD_Y, CORE_HEAD_W, CORE_HEAD_H, 'CORE FUNCTIONALITY')
    core_body = slide.shapes.add_textbox(CORE_BODY_X, CORE_BODY_Y, CORE_BODY_W, CORE_BODY_H)
    set_bullet_lines(core_body, bullets, size=14)

    return image_added


def main() -> int:
    if not PPTX_PATH.exists():
        print(f'Missing deck: {PPTX_PATH}')
        return 1

    if not SCREENSHOT_DIR.exists():
        print(f'Missing screenshot directory: {SCREENSHOT_DIR}')
        return 1

    prs = Presentation(PPTX_PATH)

    applied = 0
    missing_files: list[str] = []

    for slide_def in SLIDE_DEFINITIONS:
        image_name = str(slide_def['image'])
        screenshot_path = SCREENSHOT_DIR / image_name

        image_added = upsert_feature_slide(
            prs,
            slide_number=int(slide_def['slide']),
            title=str(slide_def['title']),
            screenshot_path=screenshot_path,
            header_text=str(slide_def['header']),
            purpose_text=str(slide_def['purpose']),
            bullets=[str(item) for item in slide_def['bullets']],
        )
        if image_added:
            applied += 1
        else:
            missing_files.append(image_name)

    prs.save(PPTX_PATH)

    print(f'Updated slide images: {applied}')
    if missing_files:
        print('Missing screenshots:')
        for item in missing_files:
            print(f' - {item}')

    return 0


if __name__ == '__main__':
    raise SystemExit(main())
